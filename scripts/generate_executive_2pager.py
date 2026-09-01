#!/usr/bin/env python3
"""Generate a visual 2-page executive deck (2 slides) for customer + leadership."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "FOSSA_Remediation_Executive_2Pager.pptx"

# Professional palette
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
TEAL = RGBColor(0x0D, 0x94, 0x88)
ACCENT = RGBColor(0x1A, 0x73, 0x6E)
LIGHT = RGBColor(0xF4, 0xF7, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5A, 0x6A, 0x7A)
LINE = RGBColor(0xC5, 0xD0, 0xD8)
AMBER = RGBColor(0xB8, 0x6E, 0x00)
GREEN = RGBColor(0x1E, 0x7A, 0x4A)
RED_SOFT = RGBColor(0xC4, 0x4E, 0x52)
MANUAL_BG = RGBColor(0xF8, 0xF1, 0xF0)
AGENT_BG = RGBColor(0xEC, 0xF6, 0xF4)


def set_run(run, size=11, bold=False, color=DARK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_rect(slide, left, top, width, height, fill, line_color=None, radius=0.08):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    return sh


def add_textbox(slide, left, top, width, height, text, size=11, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return box


def add_label(slide, left, top, width, height, text, fill, text_color=WHITE, size=10, bold=True):
    sh = add_rect(slide, left, top, width, height, fill)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(6)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=text_color)
    return sh


def add_arrow(slide, x1, y1, x2, y2, color=TEAL):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(1.5)
    return conn


def add_header_bar(slide, title, subtitle):
    add_rect(slide, 0, 0, 13.333, 0.95, NAVY)
    add_textbox(slide, 0.45, 0.18, 9.5, 0.35, title, size=22, bold=True, color=WHITE)
    add_textbox(slide, 0.45, 0.52, 9.5, 0.3, subtitle, size=11, color=RGBColor(0xA8, 0xC5, 0xC0))
    add_textbox(slide, 10.2, 0.35, 2.8, 0.3, "Working POC  ·  Confidential", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def flow_box(slide, x, y, w, h, text, fill=WHITE, tc=DARK, size=9):
    add_rect(slide, x, y, w, h, fill, LINE)
    add_textbox(slide, x + 0.05, y + 0.06, w - 0.1, h - 0.1, text, size=size, color=tc, align=PP_ALIGN.CENTER)


def build_page1(slide):
    add_header_bar(
        slide,
        "FOSSA Multi-Agent Remediation",
        "One prompt → validated plan → draft PR → FOSSA verify  ·  Neuro SAN + Mistral Devstral + FOSSA + GitHub",
    )

    # --- LEFT: Manual vs Agent ---
    add_textbox(slide, 0.4, 1.05, 6.2, 0.25, "TODAY vs WITH MULTI-AGENT", size=10, bold=True, color=NAVY)

    add_rect(slide, 0.4, 1.35, 3.0, 5.55, MANUAL_BG, LINE)
    add_label(slide, 0.55, 1.45, 2.7, 0.28, "MANUAL", RED_SOFT, WHITE, 9)
    manual = [
        "Repo + CVEs",
        "Clone",
        "Research versions",
        "Edit pom/Gradle",
        "Run tests",
        "Open PR",
        "Wait FOSSA",
        "× 10–12 services",
    ]
    for i, step in enumerate(manual):
        flow_box(slide, 0.55, 1.85 + i * 0.58, 2.7, 0.48, step, WHITE, DARK, 9)
        if i < len(manual) - 1:
            add_arrow(slide, 1.9, 2.33 + i * 0.58, 1.9, 2.43 + i * 0.58, RED_SOFT)

    add_rect(slide, 3.55, 1.35, 3.0, 5.55, AGENT_BG, LINE)
    add_label(slide, 3.7, 1.45, 2.7, 0.28, "MULTI-AGENT", TEAL, WHITE, 9)
    agent = [
        "One prompt",
        "Fetch FOSSA",
        "Plan + validate",
        "Apply + test",
        "Draft PR",
        "Verify FOSSA",
        "SRE review",
    ]
    for i, step in enumerate(agent):
        flow_box(slide, 3.7, 1.85 + i * 0.72, 2.7, 0.55, step, WHITE, DARK, 9)
        if i < len(agent) - 1:
            add_arrow(slide, 5.05, 2.4 + i * 0.72, 5.05, 2.52 + i * 0.72, TEAL)

    add_textbox(slide, 0.55, 6.55, 5.8, 0.35, "Days  →  Hours  ·  Human only at merge", size=10, bold=True, color=GREEN)

    # --- RIGHT: Architecture ---
    add_textbox(slide, 6.85, 1.05, 6.0, 0.25, "ARCHITECTURE — LLM decides · Python enforces", size=10, bold=True, color=NAVY)

    layers = [
        (0.35, "USER / SRE", "NSFlow · CLI · natural language", NAVY, WHITE),
        (1.15, "ORCHESTRATION", "fossa_orchestrator → remediation_pipeline (Devstral)", TEAL, WHITE),
        (2.05, "POLICY & DATA", "Fetch FOSSA · Plan · ValidateRemediationPlan", AMBER, WHITE),
        (2.95, "EXECUTION", "Apply · Compile · Test · Git · Draft PR", ACCENT, WHITE),
        (3.85, "VERIFY", "VerifyFossaScan (GitHub FOSSA App)", GREEN, WHITE),
    ]
    y0 = 1.4
    for i, (h, title, sub, fill, tc) in enumerate(layers):
        y = y0 + i * 0.95
        add_rect(slide, 6.85, y, 6.0, h, fill, None, 0.06)
        add_textbox(slide, 7.0, y + 0.06, 2.2, 0.22, title, size=10, bold=True, color=tc)
        add_textbox(slide, 9.1, y + 0.08, 3.6, 0.22, sub, size=9, color=tc)
        if i < len(layers) - 1:
            add_arrow(slide, 9.85, y + h, 9.85, y + h + 0.12, WHITE)

    # External systems row
    add_rect(slide, 6.85, 6.15, 6.0, 0.75, LIGHT, LINE)
    add_textbox(slide, 7.0, 6.28, 5.6, 0.5, "FOSSA API  ·  GitHub  ·  Maven/Gradle  ·  config/repos.yaml", size=9, color=MUTED, align=PP_ALIGN.CENTER)

    # Customer OpenCode callout
    add_rect(slide, 6.85, 5.05, 6.0, 0.95, WHITE, TEAL)
    add_textbox(slide, 7.0, 5.15, 5.7, 0.75, "Customer path: OpenCode as coding worker\n(plan locked after validate · git/PR optional)", size=9, bold=True, color=TEAL, align=PP_ALIGN.CENTER)


def build_page2(slide):
    add_header_bar(
        slide,
        "End-to-End Workflow & Governance",
        "Pilot: payment-service (Maven) · user-service (Gradle)  ·  Spring Boot 21",
    )

    # --- Workflow pipeline ---
    add_textbox(slide, 0.4, 1.05, 12.5, 0.25, "WORKFLOW — 5 PHASES", size=10, bold=True, color=NAVY)

    phases = [
        ("1\nDISCOVER", "Load config\nFetch FOSSA CVEs", TEAL),
        ("2\nPLAN", "Context + deps\nSubmit plan", AMBER),
        ("3\nVALIDATE", "CVE coverage\nFOSSA versions", RED_SOFT),
        ("4\nEXECUTE", "Apply · Compile\nTest · Draft PR", ACCENT),
        ("5\nVERIFY", "Poll FOSSA\n0 sec vulns", GREEN),
    ]
    x = 0.4
    for i, (title, sub, color) in enumerate(phases):
        w = 2.35
        add_rect(slide, x, 1.4, w, 1.35, color, None, 0.06)
        add_textbox(slide, x + 0.1, 1.48, w - 0.2, 0.55, title, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.1, 2.05, w - 0.2, 0.6, sub, size=9, color=WHITE, align=PP_ALIGN.CENTER)
        if i < len(phases) - 1:
            sh = slide.shapes.add_shape(
                MSO_SHAPE.CHEVRON, Inches(x + w - 0.05), Inches(1.85), Inches(0.35), Inches(0.45)
            )
            sh.fill.solid()
            sh.fill.fore_color.rgb = LINE
            sh.line.fill.background()
        x += 2.5

    add_textbox(slide, 0.4, 2.85, 12.5, 0.25, "SRE stays in loop — draft PR only · no auto-merge", size=9, color=MUTED, align=PP_ALIGN.CENTER)

    # --- Governance gates diagram ---
    add_textbox(slide, 0.4, 3.15, 6.0, 0.25, "POLICY GATES", size=10, bold=True, color=NAVY)
    gates = [
        ("ValidateRemediationPlan", "Blocks wrong versions & skipped CVEs"),
        ("RunJavaTests", "No PR if tests fail"),
        ("VerifyFossaScan", "Success = 0 security vulns on branch"),
        ("CreatePullRequest", "Draft only · triggers GitHub scan"),
    ]
    for i, (name, desc) in enumerate(gates):
        y = 3.45 + i * 0.72
        add_rect(slide, 0.4, y, 6.0, 0.62, WHITE, LINE)
        add_label(slide, 0.55, y + 0.12, 2.5, 0.38, name[:22], TEAL, WHITE, 8)
        add_textbox(slide, 3.2, y + 0.18, 3.0, 0.35, desc, size=9, color=MUTED)

    # --- Sequence-style diagram ---
    add_textbox(slide, 6.85, 3.15, 6.0, 0.25, "AGENT + TOOL FLOW", size=10, bold=True, color=NAVY)
    actors = ["You", "Orchestrator", "Pipeline", "Tools", "GitHub", "FOSSA"]
    ax = [6.95, 7.85, 8.75, 9.65, 10.55, 11.45]
    for a, x in zip(actors, ax):
        add_label(slide, x, 3.45, 0.75, 0.32, a, NAVY, WHITE, 7)
        add_rect(slide, x + 0.32, 3.85, 0.02, 2.5, LINE)

    steps = [
        (0, 1, "Remediate repo"),
        (1, 2, "Delegate"),
        (2, 3, "Fetch · Plan · Apply"),
        (3, 4, "Draft PR"),
        (3, 5, "Verify scan"),
        (0, 4, "Review & merge"),
    ]
    y_step = 4.05
    for from_i, to_i, label in steps:
        x1, x2 = ax[from_i] + 0.33, ax[to_i] + 0.33
        add_arrow(slide, min(x1, x2), y_step, max(x1, x2), y_step, TEAL)
        add_textbox(slide, 6.9, y_step - 0.12, 5.8, 0.2, label, size=7, color=MUTED)
        y_step += 0.38

    # --- Outcomes strip ---
    add_rect(slide, 0.4, 6.05, 12.5, 1.15, NAVY, None, 0.04)
    outcomes = [
        ("PROVEN", "Multi-CVE fix\non real repos"),
        ("SAFE", "Policy in\nPython code"),
        ("FAST", "One prompt\nper service"),
        ("SCALE", "12+ repos\nparallel-ready"),
    ]
    ox = 0.55
    for title, body in outcomes:
        add_textbox(slide, ox, 6.18, 2.8, 0.25, title, size=10, bold=True, color=TEAL)
        add_textbox(slide, ox, 6.45, 2.8, 0.55, body, size=9, color=WHITE)
        ox += 3.1

    add_textbox(slide, 0.4, 7.15, 12.5, 0.25, "Neuro SAN · Mistral Devstral · FOSSA · GitHub  ·  NSFlow live demo available", size=8, color=MUTED, align=PP_ALIGN.CENTER)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    s1 = prs.slides.add_slide(blank)
    add_rect(s1, 0, 0, 13.333, 7.5, LIGHT)
    build_page1(s1)

    s2 = prs.slides.add_slide(blank)
    add_rect(s2, 0, 0, 13.333, 7.5, LIGHT)
    build_page2(s2)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
